"""
咫尺天涯社交平台 - 毕业设计答辩 PPT 生成脚本
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime
import os


class PresentationGenerator:
    """PPT 生成器"""

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)  # 16:9 宽屏
        self.prs.slide_height = Inches(7.5)

        # 配色方案
        self.COLORS = {
            'primary': RGBColor(0x1E, 0x88, 0xE5),      # 主色：蓝色
            'secondary': RGBColor(0x0D, 0x47, 0xA1),    # 深蓝
            'accent': RGBColor(0xFF, 0x57, 0x22),       # 强调色：橙色
            'success': RGBColor(0x4C, 0xAF, 0x50),      # 成功色：绿色
            'white': RGBColor(0xFF, 0xFF, 0xFF),
            'light_gray': RGBColor(0xF5, 0xF5, 0xF5),
            'dark_gray': RGBColor(0x42, 0x42, 0x42),
            'text': RGBColor(0x21, 0x21, 0x21),
        }

    def add_background(self, slide, color=None):
        """添加背景色"""
        if color is None:
            color = self.COLORS['white']
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_shape_background(self, slide, left, top, width, height, color, alpha=None):
        """添加形状背景"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def add_text_box(self, slide, left, top, width, height, text, font_size=18,
                     bold=False, color=None, alignment=PP_ALIGN.LEFT, font_name='微软雅黑'):
        """添加文本框"""
        if color is None:
            color = self.COLORS['text']

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = alignment
        return txBox

    def add_bullet_points(self, slide, left, top, width, height, items, font_size=16,
                          spacing=Pt(8), color=None, bullet_char='✓'):
        """添加要点列表"""
        if color is None:
            color = self.COLORS['text']

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = f"  {bullet_char} {item}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.name = '微软雅黑'
            p.space_after = spacing

        return txBox

    def create_cover_slide(self):
        """创建封面幻灯片"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # 空白布局
        self.add_background(slide, self.COLORS['primary'])

        # 添加装饰条
        self.add_shape_background(
            slide, Inches(0), Inches(0),
            Inches(13.333), Inches(0.3),
            self.COLORS['secondary']
        )
        self.add_shape_background(
            slide, Inches(0), Inches(7.2),
            Inches(13.333), Inches(0.3),
            self.COLORS['accent']
        )

        # 标题
        self.add_text_box(
            slide, Inches(1), Inches(1.8), Inches(11.333), Inches(1.5),
            '咫尺天涯社交平台',
            font_size=54, bold=True, color=self.COLORS['white'],
            alignment=PP_ALIGN.CENTER
        )

        # 副标题
        self.add_text_box(
            slide, Inches(1), Inches(3.3), Inches(11.333), Inches(0.8),
            '基于 Python Flask 的多功能社交网络系统设计与实现',
            font_size=28, color=self.COLORS['light_gray'],
            alignment=PP_ALIGN.CENTER
        )

        # 分隔线
        self.add_shape_background(
            slide, Inches(4), Inches(4.3), Inches(5.333), Inches(0.05),
            self.COLORS['accent']
        )

        # 信息
        info_items = [
            '汇报人：XXX',
            '指导老师：XXX',
            f'日期：{datetime.now().strftime("%Y年%m月%d日")}'
        ]

        y_pos = 4.7
        for item in info_items:
            self.add_text_box(
                slide, Inches(1), Inches(y_pos), Inches(11.333), Inches(0.5),
                item, font_size=22, color=self.COLORS['light_gray'],
                alignment=PP_ALIGN.CENTER
            )
            y_pos += 0.5

    def create_outline_slide(self):
        """创建目录幻灯片"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_background(slide)

        # 标题
        self.add_shape_background(
            slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2),
            self.COLORS['primary']
        )
        self.add_text_box(
            slide, Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.8),
            '目  录', font_size=40, bold=True, color=self.COLORS['white'],
            alignment=PP_ALIGN.CENTER
        )

        # 目录项
        outline_items = [
            ('01', '项目概述', '背景、目标、技术栈'),
            ('02', '系统架构', '项目结构、技术选型'),
            ('03', '数据库设计', '数据模型、关系设计'),
            ('04', '核心功能', '用户认证、动态发布、新闻系统'),
            ('05', '社交互动', '点赞、评论、关注系统'),
            ('06', '技术亮点', '操作日志、热度算法'),
            ('07', '系统演示', '界面展示、功能演示'),
            ('08', '总结展望', '项目总结、未来规划')
        ]

        # 两列布局
        for i, (num, title, desc) in enumerate(outline_items):
            col = i // 4
            row = i % 4

            x = Inches(1.5 + col * 6)
            y = Inches(1.8 + row * 1.3)

            # 编号圆圈
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                x, y, Inches(0.7), Inches(0.7)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = self.COLORS['primary']
            circle.line.fill.background()
            tf = circle.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = num
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.COLORS['white']
            p.alignment = PP_ALIGN.CENTER

            # 标题和描述
            self.add_text_box(
                slide, x + Inches(0.9), y -
                Inches(0.1), Inches(4.5), Inches(0.4),
                title, font_size=24, bold=True, color=self.COLORS['dark_gray']
            )
            self.add_text_box(
                slide, x + Inches(0.9), y +
                Inches(0.3), Inches(4.5), Inches(0.3),
                desc, font_size=16, color=RGBColor(0x75, 0x75, 0x75)
            )

    def create_section_slide(self, slide_num, section_title, icon='●'):
        """创建章节分隔幻灯片"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_background(slide, self.COLORS['primary'])

        # 大编号
        self.add_text_box(
            slide, Inches(0), Inches(2), Inches(13.333), Inches(2),
            f'{slide_num:02d}', font_size=120, bold=True,
            color=RGBColor(0xE0, 0xE0, 0xE0),
            alignment=PP_ALIGN.CENTER
        )

        # 章节标题
        self.add_text_box(
            slide, Inches(1), Inches(3), Inches(11.333), Inches(1),
            f'{icon} {section_title}', font_size=48, bold=True,
            color=self.COLORS['white'], alignment=PP_ALIGN.CENTER
        )

    def create_content_slide(self, title, content_type='text', items=None,
                             code_snippet=None, image_path=None, subtitle=None):
        """创建内容幻灯片"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_background(slide)

        # 标题栏
        self.add_shape_background(
            slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2),
            self.COLORS['primary']
        )
        self.add_text_box(
            slide, Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.7),
            title, font_size=36, bold=True, color=self.COLORS['white'],
            alignment=PP_ALIGN.CENTER
        )

        # 副标题
        if subtitle:
            self.add_text_box(
                slide, Inches(0.5), Inches(1.3), Inches(12.333), Inches(0.5),
                subtitle, font_size=20, color=RGBColor(0x75, 0x75, 0x75),
                alignment=PP_ALIGN.CENTER
            )
            content_top = Inches(2.0)
        else:
            content_top = Inches(1.5)

        # 内容区域
        if content_type == 'text' and items:
            self.add_bullet_points(
                slide, Inches(1), content_top, Inches(11.333), Inches(5),
                items, font_size=20, spacing=Pt(12)
            )

        elif content_type == 'code' and code_snippet:
            # 代码框
            code_box = self.add_shape_background(
                slide, Inches(1), content_top, Inches(11.333), Inches(5),
                self.COLORS['light_gray']
            )

            txBox = slide.shapes.add_textbox(
                Inches(1.2), content_top +
                Inches(0.2), Inches(10.933), Inches(4.6)
            )
            tf = txBox.text_frame
            tf.word_wrap = True

            for i, line in enumerate(code_snippet.split('\n')):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(14)
                p.font.name = 'Consolas'
                p.font.color.rgb = self.COLORS['dark_gray']

        elif content_type == 'image' and image_path:
            if os.path.exists(image_path):
                slide.shapes.add_picture(
                    image_path,
                    Inches(2), content_top,
                    Inches(9.333), Inches(4.5)
                )

    def create_two_column_slide(self, title, left_title, left_items,
                                right_title, right_items):
        """创建双列布局幻灯片"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_background(slide)

        # 标题栏
        self.add_shape_background(
            slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2),
            self.COLORS['primary']
        )
        self.add_text_box(
            slide, Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.7),
            title, font_size=36, bold=True, color=self.COLORS['white'],
            alignment=PP_ALIGN.CENTER
        )

        # 左列
        self.add_shape_background(
            slide, Inches(0.5), Inches(1.5), Inches(6), Inches(0.6),
            self.COLORS['primary']
        )
        self.add_text_box(
            slide, Inches(0.5), Inches(1.5), Inches(6), Inches(0.6),
            left_title, font_size=24, bold=True, color=self.COLORS['white'],
            alignment=PP_ALIGN.CENTER
        )

        self.add_bullet_points(
            slide, Inches(0.7), Inches(2.3), Inches(5.6), Inches(4.5),
            left_items, font_size=18
        )

        # 右列
        self.add_shape_background(
            slide, Inches(6.833), Inches(1.5), Inches(6), Inches(0.6),
            self.COLORS['accent']
        )
        self.add_text_box(
            slide, Inches(6.833), Inches(1.5), Inches(6), Inches(0.6),
            right_title, font_size=24, bold=True, color=self.COLORS['white'],
            alignment=PP_ALIGN.CENTER
        )

        self.add_bullet_points(
            slide, Inches(7.033), Inches(2.3), Inches(5.6), Inches(4.5),
            right_items, font_size=18
        )

    def create_summary_slide(self):
        """创建总结幻灯片"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.add_background(slide)

        # 标题
        self.add_shape_background(
            slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2),
            self.COLORS['primary']
        )
        self.add_text_box(
            slide, Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.7),
            '项目总结与展望', font_size=36, bold=True, color=self.COLORS['white'],
            alignment=PP_ALIGN.CENTER
        )

        # 左侧：项目成果
        self.add_shape_background(
            slide, Inches(0.5), Inches(1.5), Inches(6), Inches(0.6),
            self.COLORS['success']
        )
        self.add_text_box(
            slide, Inches(0.5), Inches(1.5), Inches(6), Inches(0.6),
            '项目成果', font_size=24, bold=True, color=self.COLORS['white'],
            alignment=PP_ALIGN.CENTER
        )

        achievements = [
            '完整的社交平台功能实现',
            '真实新闻抓取与展示系统',
            '完善的操作日志与审计功能',
            '响应式界面设计',
            '数据安全与用户隐私保护'
        ]
        self.add_bullet_points(
            slide, Inches(0.7), Inches(2.3), Inches(5.6), Inches(3),
            achievements, font_size=18
        )

        # 右侧：未来展望
        self.add_shape_background(
            slide, Inches(6.833), Inches(1.5), Inches(6), Inches(0.6),
            self.COLORS['accent']
        )
        self.add_text_box(
            slide, Inches(6.833), Inches(1.5), Inches(6), Inches(0.6),
            '未来展望', font_size=24, bold=True, color=self.COLORS['white'],
            alignment=PP_ALIGN.CENTER
        )

        future_plans = [
            '引入 Redis 缓存优化性能',
            '开发移动端 App',
            '集成 AI 智能推荐',
            '支持直播与群组功能',
            '大数据分析与可视化'
        ]
        self.add_bullet_points(
            slide, Inches(7.033), Inches(2.3), Inches(5.6), Inches(3),
            future_plans, font_size=18
        )

        # 底部感谢
        self.add_shape_background(
            slide, Inches(0), Inches(6.5), Inches(13.333), Inches(1),
            self.COLORS['primary']
        )
        self.add_text_box(
            slide, Inches(0), Inches(6.6), Inches(13.333), Inches(0.8),
            '感谢各位老师聆听与指导！', font_size=32, bold=True,
            color=self.COLORS['white'], alignment=PP_ALIGN.CENTER
        )

    def generate(self):
        """生成完整的 PPT"""
        print('开始生成 PPT...')

        # 1. 封面
        print('创建封面...')
        self.create_cover_slide()

        # 2. 目录
        print('创建目录...')
        self.create_outline_slide()

        # 3. 项目概述
        print('创建项目概述...')
        self.create_section_slide(1, '项目概述')

        self.create_content_slide(
            '项目背景与目标',
            subtitle='咫尺天涯 - 让距离不再是障碍',
            content_type='text',
            items=[
                '数字化时代社交需求日益增长',
                '现有平台功能单一，缺乏综合性解决方案',
                '构建多功能社交平台：动态、视频、新闻、互动',
                '实现用户之间的无缝连接与内容分享',
                '提供完善的内容管理与审核机制',
                '建立真实可信的社交环境'
            ]
        )

        self.create_two_column_slide(
            '技术栈选择',
            '后端技术',
            [
                'Python 3.12+ 编程语言',
                'Flask Web 框架',
                'SQLAlchemy ORM',
                'SQLite 数据库',
                'Flask-Login 认证',
                'Flask-WTF 表单'
            ],
            '前端技术',
            [
                'HTML5 + CSS3 + JavaScript',
                'Bootstrap 5 响应式',
                'Font Awesome 图标',
                'jQuery + AJAX',
                '异步交互体验',
                '移动端适配'
            ]
        )

        # 4. 系统架构
        print('创建系统架构...')
        self.create_section_slide(2, '系统架构设计')

        self.create_content_slide(
            '项目结构设计',
            subtitle='清晰的分层架构，便于维护与扩展',
            content_type='text',
            items=[
                'app.py - 主应用入口，路由注册',
                'models/ - 数据库模型定义（用户、帖子、新闻等）',
                'routes/ - 路由模块（新闻、删除等）',
                'forms/ - 表单验证与处理',
                'templates/ - HTML 模板（20+ 页面）',
                'static/ - 静态资源（CSS、JS、上传文件）',
                '定时任务 - 每日自动新闻抓取'
            ]
        )

        self.create_content_slide(
            '系统架构分层',
            subtitle='MVC 架构模式',
            content_type='text',
            items=[
                '表示层：HTML/CSS/JS + Bootstrap，用户交互界面',
                '业务层：Flask Routes，处理业务逻辑',
                '数据层：SQLAlchemy ORM，数据库操作封装',
                '存储层：SQLite，持久化数据存储',
                '用户认证、内容管理、社交互动、消息推送各司其职',
                '模块化设计，高内聚低耦合'
            ]
        )

        # 5. 数据库设计
        print('创建数据库设计...')
        self.create_section_slide(3, '数据库设计')

        self.create_two_column_slide(
            '核心数据模型',
            '用户与内容',
            [
                'User - 用户信息、认证、关系',
                'Post - 动态帖子（图文视频）',
                'Video - 视频内容管理',
                'News - 新闻资讯（真实数据）',
                'Comment - 评论系统',
                'Message - 私信消息'
            ],
            '互动与日志',
            [
                'PostLike - 帖子点赞',
                'VideoLike - 视频点赞',
                'NewsLike - 新闻点赞',
                'Follows - 关注关系',
                'NewsOperationLog - 操作日志',
                'ContentHistory - 内容历史'
            ]
        )

        self.create_content_slide(
            '数据库设计亮点',
            subtitle='规范设计，保障数据完整性',
            content_type='text',
            items=[
                '外键约束确保数据关联完整性',
                '唯一约束防止重复数据（用户名、邮箱）',
                '默认值设置提升用户体验',
                '时间戳自动记录创建与更新时间',
                '操作日志表支持审计追溯',
                '索引优化查询性能'
            ]
        )

        # 6. 核心功能
        print('创建核心功能...')
        self.create_section_slide(4, '核心功能模块')

        self.create_two_column_slide(
            '用户认证系统',
            '安全机制',
            [
                '密码加密存储（Werkzeug）',
                '图形验证码防刷',
                '会话管理（Flask-Login）',
                '记住登录状态',
                '登录状态实时同步',
                '安全退出机制'
            ],
            '功能特性',
            [
                '用户注册（头像选择/上传）',
                '用户登录（用户名+密码）',
                '个人资料编辑',
                '头像更换',
                '在线状态管理',
                '权限控制（登录后操作）'
            ]
        )

        self.create_two_column_slide(
            '动态发布系统',
            '发布功能',
            [
                '图文动态发布',
                '视频动态上传',
                '实时预览',
                '文件大小验证',
                '格式限制检查',
                '自动缩略图生成'
            ],
            '管理功能',
            [
                '编辑已发布动态',
                '删除违规内容',
                '点赞与取消',
                '评论互动',
                '热度计算排行',
                '内容修改历史'
            ]
        )

        self.create_content_slide(
            '新闻系统 - 真实数据抓取',
            subtitle='集成多个权威新闻源，提供真实资讯',
            content_type='text',
            items=[
                '新闻源：新华网、央视新闻、澎湃新闻等',
                '定时任务：每日 08:00 自动抓取',
                '数据去重：避免重复新闻',
                '日历视图：按月浏览新闻',
                '按日期查看：精准定位',
                '关键字搜索：标题+内容+摘要',
                '来源筛选：快速过滤'
            ]
        )

        # 7. 社交互动
        print('创建社交互动...')
        self.create_section_slide(5, '社交互动系统')

        self.create_content_slide(
            '点赞与评论系统',
            subtitle='实时互动，增强用户参与感',
            content_type='text',
            items=[
                'AJAX 无刷新点赞/取消',
                '实时更新点赞数显示',
                '防止重复点赞（唯一约束）',
                '评论发布与展示',
                '评论数自动统计',
                '操作日志完整记录'
            ]
        )

        self.create_content_slide(
            '关注与粉丝系统',
            subtitle='建立社交关系网络',
            content_type='text',
            items=[
                '关注/取消关注功能',
                '关注列表与粉丝列表',
                '关注数与粉丝数统计',
                '动态推送（关注的人的内容）',
                '私信系统（一对一沟通）',
                '在线状态实时显示'
            ]
        )

        # 8. 技术亮点
        print('创建技术亮点...')
        self.create_section_slide(6, '技术创新与亮点')

        self.create_content_slide(
            '操作日志系统 - 完整审计追踪',
            subtitle='每次操作都有记录，数据变更可追溯',
            content_type='text',
            items=[
                '记录查看操作：浏览数、时间、IP',
                '记录点赞操作：前后值对比',
                '记录评论操作：评论数变化',
                '记录设备信息：User-Agent',
                '支持安全审计与数据分析',
                '防止数据篡改，确保真实性'
            ]
        )

        self.create_content_slide(
            '热度排行算法 - 智能排序',
            subtitle='多维度加权计算，时间衰减因子',
            content_type='text',
            items=[
                '基础评分 = 点赞×3 + 浏览×0.5 + 评论×2',
                '时间衰减因子 = 24 / (年龄小时 + 1)',
                '最终得分 = 基础评分 × 衰减因子',
                '每日/每周/每月热度排行',
                '新内容获得更高曝光机会',
                '鼓励优质内容创作'
            ]
        )

        self.create_content_slide(
            '实时数据更新 - 用户体验优化',
            subtitle='前后端联动，无刷新操作',
            content_type='text',
            items=[
                'Fetch API 实现异步请求',
                'JavaScript 动态更新 DOM',
                '点赞状态即时反馈',
                '数字实时变化',
                '减少页面刷新',
                '提升交互流畅度'
            ]
        )

        # 9. 系统演示
        print('创建系统演示...')
        self.create_section_slide(7, '系统功能演示')

        self.create_content_slide(
            '功能展示',
            subtitle='完整的功能实现，流畅的用户体验',
            content_type='text',
            items=[
                '首页：动态列表、发布框、导航栏',
                '新闻页面：列表、搜索、日历视图',
                '用户资料：个人信息、发布历史、关注列表',
                '视频播放：在线播放、点赞评论',
                '私信聊天：实时消息、附件发送',
                '热度排行：每日/每周/每月排行',
                '响应式设计：适配各种设备'
            ]
        )

        # 10. 测试
        print('创建测试...')
        self.create_section_slide(8, '测试与优化')

        self.create_two_column_slide(
            '功能测试与性能',
            '测试结果',
            [
                '用户认证：✅ 通过',
                '动态发布：✅ 通过',
                '新闻抓取：✅ 通过',
                '点赞评论：✅ 通过',
                '操作日志：✅ 通过',
                '搜索功能：✅ 通过'
            ],
            '性能优化',
            [
                '首页加载 < 2 秒',
                '新闻搜索 < 3 秒',
                '点赞响应 < 0.5 秒',
                '数据库索引优化',
                'SQL 查询优化',
                '静态资源缓存'
            ]
        )

        # 11. 总结
        print('创建总结...')
        self.create_summary_slide()

        # 保存 PPT
        output_path = os.path.join(
            os.path.dirname(os.path.abspath(__file__)),
            '咫尺天涯社交平台_毕业设计答辩.pptx'
        )
        self.prs.save(output_path)
        print(f'\n✅ PPT 生成成功！')
        print(f'📁 文件位置：{output_path}')
        print(f'📊 幻灯片数量：{len(self.prs.slides)} 页')


if __name__ == '__main__':
    generator = PresentationGenerator()
    generator.generate()
